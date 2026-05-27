"""
Generates PanelIQ_Executive_Deck_May2026.pptx — 5 slides for top management
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as RGB
from pptx.enum.text import PP_ALIGN
import os

BASE = r'C:\Users\KarthikShanmugam\ClaudePOC\paneliq'

# ── Colours ───────────────────────────────────────────────────────────────────
DARK_BLUE = RGB(0x1a, 0x3a, 0x5c)
TEAL      = RGB(0x00, 0x99, 0xcc)
GREEN     = RGB(0x00, 0xa8, 0x6b)
WHITE     = RGB(0xff, 0xff, 0xff)
GRAY      = RGB(0x2d, 0x37, 0x48)
LIGHT     = RGB(0xf0, 0xf7, 0xfb)
ORANGE    = RGB(0xe6, 0x7e, 0x22)
MID_GRAY  = RGB(0xa0, 0xb8, 0xd0)
RED_SOFT  = RGB(0xc0, 0x39, 0x2b)
PALE_GREEN= RGB(0xe8, 0xf8, 0xf0)
PALE_BLUE = RGB(0xe8, 0xf4, 0xfb)


def prs_new():
    p = Presentation()
    p.slide_width  = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = lw
    else:
        s.line.fill.background()
    return s

def circle(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
       italic=False, wrap=True):
    s = slide.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color if color else GRAY
    return s

def two_run_tb(slide, l, t, w, h,
               text1, size1, bold1, color1,
               text2, size2, bold2, color2,
               align=PP_ALIGN.LEFT):
    """Single textbox with two inline runs — bold label + normal body."""
    s = slide.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r1 = p.add_run(); r1.text = text1
    r1.font.size = Pt(size1); r1.font.bold = bold1; r1.font.color.rgb = color1
    r2 = p.add_run(); r2.text = text2
    r2.font.size = Pt(size2); r2.font.bold = bold2; r2.font.color.rgb = color2
    return s

def stat_block(slide, l, t, w, h, big, label, bg, big_color, label_color):
    """Big number + label stacked inside a coloured card."""
    box(slide, l, t, w, h, bg, TEAL, Pt(1))
    tb(slide, big,   l, t + Inches(0.25), w, Inches(0.9),
       size=42, bold=True, color=big_color, align=PP_ALIGN.CENTER)
    tb(slide, label, l, t + Inches(1.05), w, Inches(0.8),
       size=11, bold=False, color=label_color, align=PP_ALIGN.CENTER, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — HOOK
# ══════════════════════════════════════════════════════════════════════════════
def slide1_hook(prs):
    s = blank(prs)

    # Full dark background
    box(s, 0, 0, prs.slide_width, prs.slide_height, DARK_BLUE)

    # Teal left accent bar
    box(s, 0, 0, Inches(0.18), prs.slide_height, TEAL)

    # Top-right pill: POC status
    box(s, Inches(9.8), Inches(0.35), Inches(3.1), Inches(0.6),
        GREEN, line=None)
    tb(s, '✓  POC Built — Phase 1 is 80% Complete',
       Inches(9.85), Inches(0.38), Inches(3.0), Inches(0.5),
       size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Main headline
    tb(s, 'What if your leadership team',
       Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.9),
       size=36, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)
    tb(s, 'could get data answers instantly —',
       Inches(0.5), Inches(1.95), Inches(12.5), Inches(0.9),
       size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    tb(s, 'instead of waiting 1–3 days?',
       Inches(0.5), Inches(2.75), Inches(12.5), Inches(0.9),
       size=40, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

    # Divider line
    box(s, Inches(0.5), Inches(3.85), Inches(4.5), Pt(2), TEAL)

    # Sub line
    tb(s, 'PanelIQ is an AI-powered analytics assistant that gives Borderless Access '
          'leadership instant, self-service access to panel performance data — '
          'in plain English, with no technical expertise required.',
       Inches(0.5), Inches(4.05), Inches(10.5), Inches(1.2),
       size=14, color=MID_GRAY, wrap=True)

    # Bottom branding
    tb(s, 'PanelIQ  |  Borderless Access  |  May 2026',
       Inches(0.5), Inches(6.9), Inches(8), Inches(0.45),
       size=11, color=RGB(0x55, 0x70, 0x90))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
def slide2_problem(prs):
    s = blank(prs)
    box(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)

    # Header bar
    box(s, 0, 0, prs.slide_width, Inches(0.95), DARK_BLUE)
    box(s, 0, Inches(0.95), prs.slide_width, Pt(4), TEAL)
    tb(s, 'The Challenge — How It Works Today',
       Inches(0.4), Inches(0.12), Inches(12), Inches(0.7),
       size=22, bold=True, color=WHITE)

    # ── LEFT: TODAY (pain) ──────────────────────────────────────────────────
    box(s, Inches(0.3), Inches(1.1), Inches(5.9), Inches(5.9),
        RGB(0xff, 0xf5, 0xf5), RED_SOFT, Pt(1))

    tb(s, '😟  Today', Inches(0.5), Inches(1.25), Inches(5.5), Inches(0.55),
       size=16, bold=True, color=RED_SOFT)

    pain_items = [
        ('Static dashboards only',
         'Tableau & SSRS show pre-defined views. Anything dynamic requires a manual request.'),
        ('1–3 business days',
         'Every ad-hoc question is routed to Technology, Finance, or Analytics teams.'),
        ('Once a month or quarter',
         'Leadership receives dynamic insights only during scheduled business reviews.'),
        ('No self-service',
         'Non-technical users have no independent way to access live or dynamic data.'),
    ]
    for i, (bold_t, rest) in enumerate(pain_items):
        two_run_tb(s,
                   Inches(0.5), Inches(1.95) + i * Inches(1.18),
                   Inches(5.5), Inches(1.05),
                   '✗  ' + bold_t + '\n', 11, True, RED_SOFT,
                   rest, 10.5, False, GRAY)

    # ── RIGHT: WITH PANELIQ (benefit) ───────────────────────────────────────
    box(s, Inches(6.6), Inches(1.1), Inches(6.4), Inches(5.9),
        PALE_GREEN, GREEN, Pt(1))

    tb(s, '😊  With PanelIQ', Inches(6.8), Inches(1.25), Inches(6.0), Inches(0.55),
       size=16, bold=True, color=GREEN)

    good_items = [
        ('Ask in plain English',
         'Type any question — get a data table, chart, and analyst summary in seconds.'),
        ('Seconds, not days',
         'No waiting. No tickets. No dependency on other teams.'),
        ('Always available',
         'Leadership can query data any time — daily if needed, not just in reviews.'),
        ('Secure self-service',
         'Role-based access ensures each user sees only what they are permitted to see.'),
    ]
    for i, (bold_t, rest) in enumerate(good_items):
        two_run_tb(s,
                   Inches(6.8), Inches(1.95) + i * Inches(1.18),
                   Inches(5.9), Inches(1.05),
                   '✓  ' + bold_t + '\n', 11, True, GREEN,
                   rest, 10.5, False, GRAY)

    # Centre arrow
    tb(s, '→', Inches(6.05), Inches(3.65), Inches(0.55), Inches(0.7),
       size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — HOW IT WORKS + WHAT YOU CAN ASK
# ══════════════════════════════════════════════════════════════════════════════
def slide3_solution(prs):
    s = blank(prs)
    box(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)

    box(s, 0, 0, prs.slide_width, Inches(0.95), DARK_BLUE)
    box(s, 0, Inches(0.95), prs.slide_width, Pt(4), TEAL)
    tb(s, 'PanelIQ — How It Works',
       Inches(0.4), Inches(0.12), Inches(12), Inches(0.7),
       size=22, bold=True, color=WHITE)

    # ── 4 step boxes ────────────────────────────────────────────────────────
    steps = [
        ('1', 'Ask',        'Type a question in plain English',                    TEAL),
        ('2', 'Translate',  'Claude AI converts it to a secure SQL query',         DARK_BLUE),
        ('3', 'Query',      'Runs against the internal KpiReports database',       TEAL),
        ('4', 'Answer',     'Chart + data table + analyst summary in seconds',     GREEN),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        lft = Inches(0.35) + i * Inches(3.25)
        tp  = Inches(1.1)

        box(s, lft, tp, Inches(3.0), Inches(2.9), LIGHT, color, Pt(1))

        c = circle(s, lft + Inches(1.0), tp + Inches(0.2), Inches(1.0), Inches(1.0), color)
        ctf = c.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(22); cr.font.bold = True; cr.font.color.rgb = WHITE

        tb(s, title, lft + Inches(0.1), tp + Inches(1.35),
           Inches(2.8), Inches(0.45), size=13, bold=True, color=color,
           align=PP_ALIGN.CENTER)
        tb(s, desc, lft + Inches(0.1), tp + Inches(1.88),
           Inches(2.8), Inches(0.9), size=10.5, color=GRAY,
           align=PP_ALIGN.CENTER, wrap=True)

        if i < 3:
            tb(s, '→', Inches(3.3) + i * Inches(3.25), Inches(2.25),
               Inches(0.35), Inches(0.5), size=22, bold=True, color=TEAL)

    # ── Sample questions panel ───────────────────────────────────────────────
    box(s, Inches(0.35), Inches(4.15), Inches(12.6), Inches(2.85),
        DARK_BLUE)
    tb(s, 'What your team can ask — examples',
       Inches(0.55), Inches(4.28), Inches(12.0), Inches(0.4),
       size=12, bold=True, color=TEAL)

    questions = [
        '"Show me completes by market for Q1 2026"',
        '"Which vendors have the highest screenout rate this year?"',
        '"What is the click rate and conversion rate by country?"',
        '"Compare monthly completes trend for the last 12 months"',
        '"Which projects are live right now?"',
        '"Give me overall panel health summary for this year"',
    ]
    for i, q in enumerate(questions):
        col = i % 3; row = i // 3
        lft = Inches(0.55) + col * Inches(4.15)
        tp  = Inches(4.78) + row * Inches(0.95)
        box(s, lft, tp, Inches(3.95), Inches(0.8),
            RGB(0x22, 0x44, 0x66), TEAL, Pt(1))
        tb(s, q, lft + Inches(0.12), tp + Inches(0.12),
           Inches(3.7), Inches(0.6), size=10, color=WHITE, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — INVESTMENT & STATUS
# ══════════════════════════════════════════════════════════════════════════════
def slide4_investment(prs):
    s = blank(prs)
    box(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)

    box(s, 0, 0, prs.slide_width, Inches(0.95), DARK_BLUE)
    box(s, 0, Inches(0.95), prs.slide_width, Pt(4), TEAL)
    tb(s, 'Investment & Implementation',
       Inches(0.4), Inches(0.12), Inches(12), Inches(0.7),
       size=22, bold=True, color=WHITE)

    # ── 3 stat blocks (top row) ──────────────────────────────────────────────
    stats = [
        ('~$50/mo',   'Maximum ongoing\noperating cost\n(API only)', PALE_BLUE,  DARK_BLUE, GRAY),
        ('80%',       'Phase 1 already built\n(POC validated —\nfaster to go live)', PALE_GREEN, GREEN,     GRAY),
        ('Seconds',   'Time to answer\nany panel question\nvs. 1–3 days today',   LIGHT,     TEAL,      GRAY),
    ]
    for i, (big, label, bg, big_c, lbl_c) in enumerate(stats):
        lft = Inches(0.4) + i * Inches(4.3)
        stat_block(s, lft, Inches(1.1), Inches(4.0), Inches(2.1),
                   big, label, bg, big_c, lbl_c)

    # ── Scenario comparison ──────────────────────────────────────────────────
    tb(s, 'Implementation Options', Inches(0.4), Inches(3.4), Inches(12), Inches(0.4),
       size=13, bold=True, color=DARK_BLUE)

    scenarios = [
        ('Scenario 1', 'Technology Team',        '~6–8 months',  '~$1,200–$1,600\nlicence cost', LIGHT,     DARK_BLUE),
        ('Scenario 2', 'Analytics Team',         '~12–15 months','~$2,400–$4,500\nlicence cost', LIGHT,     DARK_BLUE),
        ('Scenario 3', 'Hybrid  ★ Recommended',  '~8–10 months', '~$1,600–$2,000\nlicence cost', PALE_GREEN, GREEN),
    ]
    for i, (sc, team, timeline, cost, bg, color) in enumerate(scenarios):
        lft = Inches(0.4) + i * Inches(4.3)
        tp  = Inches(3.85)
        box(s, lft, tp, Inches(4.0), Inches(3.0), bg, color, Pt(1))

        tb(s, sc,       lft + Inches(0.15), tp + Inches(0.15),
           Inches(3.7), Inches(0.35), size=10, bold=True, color=MID_GRAY)
        tb(s, team,     lft + Inches(0.15), tp + Inches(0.48),
           Inches(3.7), Inches(0.45), size=13, bold=True, color=color)
        tb(s, f'⏱  Timeline:  {timeline}',
           lft + Inches(0.15), tp + Inches(1.05),
           Inches(3.7), Inches(0.45), size=11, color=GRAY)
        tb(s, f'💰  Dev licences:  {cost}',
           lft + Inches(0.15), tp + Inches(1.55),
           Inches(3.7), Inches(0.7), size=11, color=GRAY, wrap=True)
        tb(s, 'Ongoing API: ~$35–$50/month  (all scenarios)',
           lft + Inches(0.15), tp + Inches(2.55),
           Inches(3.7), Inches(0.35), size=9.5, italic=True, color=color)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — THE ASK
# ══════════════════════════════════════════════════════════════════════════════
def slide5_ask(prs):
    s = blank(prs)

    # Dark background
    box(s, 0, 0, prs.slide_width, prs.slide_height, DARK_BLUE)
    box(s, 0, Inches(0), prs.slide_width, Pt(4), TEAL)

    # Headline
    tb(s, 'One Decision. Three Steps. Live in Weeks.',
       Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8),
       size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── 3 action blocks ─────────────────────────────────────────────────────
    actions = [
        ('01', 'Approve',
         'Confirm scope and\nselect the team model\n(Hybrid recommended)',
         TEAL,      WHITE),
        ('02', 'Assign',
         'Name Analytics leads\n+ Technology support\nProcure Claude licences',
         DARK_BLUE, TEAL),
        ('03', 'Launch',
         'Phase 1 goes live at\npaneliq.borderlessaccess.com\nwithin 5–7 weeks',
         GREEN,     WHITE),
    ]
    for i, (num, title, desc, bg, tc) in enumerate(actions):
        lft = Inches(0.8) + i * Inches(3.95)
        tp  = Inches(1.5)

        box(s, lft, tp, Inches(3.65), Inches(4.3), bg)

        c = circle(s, lft + Inches(1.3), tp + Inches(0.25),
                   Inches(1.05), Inches(1.05), WHITE)
        ctf = c.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(20); cr.font.bold = True; cr.font.color.rgb = bg

        tb(s, title, lft + Inches(0.1), tp + Inches(1.5),
           Inches(3.45), Inches(0.6), size=18, bold=True, color=WHITE,
           align=PP_ALIGN.CENTER)
        tb(s, desc, lft + Inches(0.15), tp + Inches(2.25),
           Inches(3.35), Inches(1.8), size=12, color=WHITE,
           align=PP_ALIGN.CENTER, wrap=True)

        if i < 2:
            tb(s, '→', Inches(4.4) + i * Inches(3.95), Inches(3.3),
               Inches(0.35), Inches(0.55), size=24, bold=True, color=TEAL)

    # Bottom line
    box(s, Inches(1.5), Inches(6.2), Inches(10.3), Pt(2), TEAL)
    tb(s, 'paneliq.borderlessaccess.com  |  Borderless Access  |  May 2026',
       Inches(1.5), Inches(6.35), Inches(10.3), Inches(0.5),
       size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── MAIN ──────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    prs = prs_new()
    slide1_hook(prs)
    slide2_problem(prs)
    slide3_solution(prs)
    slide4_investment(prs)
    slide5_ask(prs)

    out = os.path.join(BASE, 'PanelIQ_Executive_Deck_May2026.pptx')
    prs.save(out)
    print(f'Saved: {out}')
