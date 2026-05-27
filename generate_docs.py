"""
Generates PanelIQ_Proposal_May2026.docx and PanelIQ_Proposal_May2026.pptx
"""
from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGB, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PPTRGB
from pptx.enum.text import PP_ALIGN

import os

BASE = r'C:\Users\KarthikShanmugam\ClaudePOC\paneliq'

# ── Shared colour constants ────────────────────────────────────────────────────
D_DARK_BLUE  = DocxRGB(0x1a, 0x3a, 0x5c)
D_TEAL       = DocxRGB(0x00, 0x99, 0xcc)
D_WHITE      = DocxRGB(0xff, 0xff, 0xff)
D_GRAY       = DocxRGB(0x2d, 0x37, 0x48)
D_LIGHT_ROW  = DocxRGB(0xf0, 0xf7, 0xfb)
D_GREEN      = DocxRGB(0x00, 0xa8, 0x6b)

P_DARK_BLUE  = PPTRGB(0x1a, 0x3a, 0x5c)
P_TEAL       = PPTRGB(0x00, 0x99, 0xcc)
P_WHITE      = PPTRGB(0xff, 0xff, 0xff)
P_GRAY       = PPTRGB(0x2d, 0x37, 0x48)
P_LIGHT      = PPTRGB(0xf0, 0xf7, 0xfb)
P_GREEN      = PPTRGB(0x00, 0xa8, 0x6b)
P_PURPLE     = PPTRGB(0x6c, 0x5c, 0xe7)
P_MID_GRAY   = PPTRGB(0xa0, 0xb8, 0xd0)


# ══════════════════════════════════════════════════════════════════════════════
#  WORD DOCUMENT
# ══════════════════════════════════════════════════════════════════════════════

def cell_bg(cell, hex_str):
    tc = cell._tc
    pr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'),   'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'),  hex_str)
    pr.append(shd)

def h1(doc, text):
    p = doc.add_heading(text, level=1)
    for r in p.runs:
        r.font.color.rgb = D_DARK_BLUE
    return p

def h2(doc, text):
    p = doc.add_heading(text, level=2)
    for r in p.runs:
        r.font.color.rgb = D_TEAL
    return p

def bullet(doc, bold_part, rest=''):
    p = doc.add_paragraph(style='List Bullet')
    if bold_part:
        r = p.add_run(bold_part + (': ' if rest else ''))
        r.bold = True
        r.font.color.rgb = D_DARK_BLUE
        r.font.size = DocxPt(10.5)
    if rest:
        r2 = p.add_run(rest)
        r2.font.size = DocxPt(10.5)

def doc_table(doc, headers, rows, highlight_last=False):
    t = doc.add_table(rows=len(rows) + 1, cols=len(headers))
    t.style = 'Table Grid'
    # Header row
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]
        cell_bg(c, '1a3a5c')
        r = c.paragraphs[0].add_run(h)
        r.bold = True
        r.font.color.rgb = D_WHITE
        r.font.size = DocxPt(10)
    # Data rows
    for ri, row in enumerate(rows):
        is_total = highlight_last and ri == len(rows) - 1
        is_first_col_header = (len(headers) > 1 and ri >= 0)
        for ci, val in enumerate(row):
            c = t.rows[ri + 1].cells[ci]
            if is_total:
                cell_bg(c, '1a3a5c')
                r = c.paragraphs[0].add_run(val)
                r.bold = True
                r.font.color.rgb = D_WHITE
                r.font.size = DocxPt(10)
            elif ci == 0 and len(headers) >= 4:
                cell_bg(c, '1a3a5c')
                r = c.paragraphs[0].add_run(val)
                r.bold = True
                r.font.color.rgb = D_WHITE
                r.font.size = DocxPt(10)
            else:
                if ri % 2 == 0:
                    cell_bg(c, 'f0f7fb')
                r = c.paragraphs[0].add_run(val)
                r.font.size = DocxPt(10)
    return t

def create_word(path):
    doc = Document()
    for s in doc.sections:
        s.top_margin    = DocxInches(1)
        s.bottom_margin = DocxInches(1)
        s.left_margin   = DocxInches(1.2)
        s.right_margin  = DocxInches(1.2)

    # ── TITLE PAGE ──────────────────────────────────────────────────────────
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = DocxPt(100)
    r = p.add_run('PanelIQ')
    r.bold = True; r.font.size = DocxPt(38); r.font.color.rgb = D_DARK_BLUE

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('AI-Powered Panel Intelligence Platform')
    r.font.size = DocxPt(20); r.font.color.rgb = D_TEAL

    for _ in range(3): doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('Proposal & Implementation Plan')
    r.font.size = DocxPt(14); r.font.color.rgb = D_GRAY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('Borderless Access  |  May 2026')
    r.font.size = DocxPt(12); r.font.color.rgb = DocxRGB(0x71, 0x8a, 0xa4)

    doc.add_page_break()

    # ── 1. EXECUTIVE SUMMARY ────────────────────────────────────────────────
    h1(doc, '1.  Executive Summary')
    doc.add_paragraph(
        'PanelIQ is an AI-powered natural language analytics platform purpose-built for Borderless Access (BA). '
        'It enables leadership and senior management to ask questions about panel performance in plain English '
        'and receive instant, data-backed answers — without relying on the Technology, Finance, or Analytics teams.'
    )
    doc.add_paragraph(
        'Today, any question beyond standard Tableau dashboards or SSRS reports takes 1–3 business days to answer. '
        'PanelIQ reduces this to seconds. The platform is built on scalable, enterprise-grade technology and '
        'is designed to grow with the business — adding new data domains, users, and capabilities over time.'
    )
    p = doc.add_paragraph()
    r = p.add_run(
        'A Proof of Concept (POC) has been developed and validated. '
        'Phase 1 of the platform is approximately 80% complete.'
    )
    r.bold = True; r.font.color.rgb = D_TEAL

    # ── 2. CURRENT STATE ────────────────────────────────────────────────────
    doc.add_paragraph()
    h1(doc, '2.  Current State & Challenges')
    h2(doc, '2.1  How Insights Are Accessed Today')
    doc.add_paragraph(
        'Leadership and senior management at BA currently rely on Tableau dashboards for standardised performance views, '
        'SSRS reports for structured operational reporting, and manual requests to internal teams for any '
        'question outside these standard reports. For dynamic or ad-hoc questions, the process involves raising '
        'a request with the Technology team, Finance team, or Analytics team, and waiting 1–3 business days '
        'for a response. Dynamic insights are typically available only during monthly or quarterly business reviews.'
    )

    h2(doc, '2.2  Key Challenges')
    challenges = [
        ('Delayed Decision-Making',   '1–3 business days for any dynamic question. Leaders cannot act on data in real time.'),
        ('Team Dependency',            'Every ad-hoc question creates additional workload for Technology, Finance, or Analytics teams.'),
        ('Limited Frequency',          'Dynamic insights are available only monthly or quarterly during business review meetings.'),
        ('Inconsistency Risk',         'Multiple teams interpreting the same metric may produce conflicting numbers.'),
        ('Scalability Bottleneck',     'As BA grows, the volume of ad-hoc data requests increases — the current process does not scale.'),
        ('No Self-Service Capability', 'Non-technical users have no independent access to dynamic or ad-hoc data.'),
    ]
    for b, d in challenges:
        bullet(doc, b, d)

    # ── 3. PROPOSED SOLUTION ────────────────────────────────────────────────
    doc.add_paragraph()
    h1(doc, '3.  Proposed Solution — PanelIQ')
    h2(doc, '3.1  What is PanelIQ?')
    doc.add_paragraph(
        'PanelIQ is a secure, web-based AI assistant accessible at paneliq.borderlessaccess.com. '
        'Users ask questions in plain English and receive answers in seconds — complete with data tables, '
        'visualisations, and AI-generated analyst summaries.'
    )

    h2(doc, '3.2  How It Works')
    steps = [
        'User types a natural language question (e.g., "Show me completes by market for Q1 2026")',
        'Claude AI translates the question into a precise, secure SQL query',
        'The query executes against the internal KpiReports database',
        'Results are returned within seconds — with a chart, data table, and analyst summary',
    ]
    for s in steps:
        p = doc.add_paragraph(style='List Number')
        p.add_run(s).font.size = DocxPt(10.5)

    h2(doc, '3.3  Access Control & Data Security')
    access = [
        ('Role-Based Access',         'Users see only data relevant to their role and business unit.'),
        ('Field-Level Restrictions',   'Sensitive fields — revenue, vendor cost, margins — are visible only to authorised users.'),
        ('Row-Level Scoping',          'A BU director sees only their BU\'s client revenue — not data from other business units.'),
        ('AI-Enforced Rules',          'Access constraints are injected into every SQL query at the AI layer and cannot be bypassed through conversation.'),
    ]
    for b, d in access:
        bullet(doc, b, d)

    h2(doc, '3.4  Features — Current & Planned')
    p = doc.add_paragraph()
    r = p.add_run('Current (POC — available now):')
    r.bold = True; r.font.color.rgb = D_DARK_BLUE

    for f in [
        'Natural language to SQL query engine',
        'Sample performance analytics — market, vendor, completes, IR%, Click Rate, Conversion Rate',
        'Automated chart type recommendation (bar, line, pie, table, metric cards)',
        'AI-generated analyst summaries with follow-up question suggestions',
        'Debounced prompt suggestions as the user types',
    ]:
        bullet(doc, '', f)

    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run('Planned (future phases):')
    r.bold = True; r.font.color.rgb = D_DARK_BLUE

    for f in [
        'Personalised question recommendations based on user role and data access',
        'Query history — previously asked questions stored and resurfaced as suggestions',
        'Vendor cost and quality analysis module',
        'Revenue by market, client, and business unit',
        'Financial module with restricted field-level access controls',
        'Predictive trend analysis and automated anomaly alerts',
    ]:
        bullet(doc, '', f)

    # ── 4. TECHNOLOGY STACK ─────────────────────────────────────────────────
    doc.add_paragraph()
    h1(doc, '4.  Technology Stack')
    doc.add_paragraph(
        'All components are industry-standard, enterprise-grade, actively maintained, '
        'and selected for long-term scalability, maintainability, and developer availability. '
        'The platform is designed to support BA for many years and can be extended with new '
        'data modules, users, and features without rebuilding the core architecture.'
    )
    doc_table(doc,
        ['Component', 'Technology', 'Purpose', 'Scalability'],
        [
            ['AI Engine',      'Anthropic Claude API',        'Natural language to SQL; summarisation; suggestions',         'API-based; no infra to manage; model improves over time'],
            ['Backend',        'Python + FastAPI',             'Query engine, business logic, access control',                'Production-grade; high concurrency; widely supported'],
            ['Frontend',       'React + Vite',                 'Chat UI, charts, data tables',                               'Industry standard; large ecosystem; easy to extend'],
            ['Analytics DB',   'SQL Server (KpiReports)',      'Source of truth for all panel data',                         'Enterprise-grade; existing internal infrastructure'],
            ['Application DB', 'SQL Server (new internal)',    'User profiles, access rules, query history, AI guidelines',   'Leverages existing SQL Server — no additional licence'],
            ['Hosting',        'Internal subdomain',           'paneliq.borderlessaccess.com',                               'No cost; scalable to 500+ concurrent users'],
            ['Version Control','GitHub',                       'Code management and team collaboration',                     'Industry standard; private repository'],
            ['Dev Assistant',  'Claude Code (Anthropic)',      'AI-assisted development — accelerates all phases',           'Significantly reduces learning curve for analytics team'],
        ]
    )

    # ── 5. PHASED ROADMAP ───────────────────────────────────────────────────
    doc.add_paragraph()
    h1(doc, '5.  Phased Roadmap')
    phases = [
        ('Phase 1 — Foundation',            '~80% complete (POC built)',          [
            'Hosted deployment at paneliq.borderlessaccess.com',
            'User authentication and login',
            'Role-based and row-level access control per business unit',
            'Sample performance analytics (markets, vendors, completes, IR%, rates)',
            'Query history and personalised question suggestions',
            'Application database setup (users, access rules, AI guidelines)',
        ]),
        ('Phase 2 — Vendor & Quality',       'Pending schema availability',         [
            'Vendor cost analysis by project, market, and period',
            'Sample quality metrics and rejection analysis',
            'Vendor quality scoring and performance comparison',
            'New SQL tables to be integrated (schema to be provided)',
        ]),
        ('Phase 3 — Commercial & Revenue',   'Pending Finance data access',          [
            'Revenue by market, client, and business unit',
            'Research cost and margin analysis',
            'Field-level financial access control',
            'Finance team data integration',
        ]),
        ('Phase 4 — Advanced Intelligence',  'Future enhancement',                  [
            'Trend forecasting and predictive insights',
            'Automated alerts for anomalies (IR% drops, vendor quality dips)',
            'Executive summary automation',
            'Integration with existing tools (Tableau, SSRS)',
        ]),
    ]
    for title, status, items in phases:
        p = doc.add_paragraph()
        r1 = p.add_run(title)
        r1.bold = True; r1.font.size = DocxPt(12); r1.font.color.rgb = D_DARK_BLUE
        r2 = p.add_run(f'   —   {status}')
        r2.italic = True; r2.font.size = DocxPt(10); r2.font.color.rgb = D_TEAL
        for item in items:
            bullet(doc, '', item)
        doc.add_paragraph()

    # ── 6. COST & TIMELINE ──────────────────────────────────────────────────
    doc.add_page_break()
    h1(doc, '6.  Cost & Timeline')
    h2(doc, '6.1  Implementation Scenarios')
    doc_table(doc,
        ['', 'Scenario 1\nTech Team', 'Scenario 2\nAnalytics Team', 'Scenario 3 ★\nHybrid (Recommended)'],
        [
            ['Team',             '1–2 Tech developers',          '2–3 Analysts + Claude AI',      'Analytics leads + Tech for infra'],
            ['Phase 1',          '4–6 weeks',                    '10–14 weeks',                    '5–7 weeks'],
            ['Phase 2',          '6–8 weeks',                    '10–14 weeks',                    '7–10 weeks'],
            ['Phase 3',          '6–8 weeks',                    '10–14 weeks',                    '7–10 weeks'],
            ['Phase 4',          '8–10 weeks',                   '14–18 weeks',                    '10–12 weeks'],
            ['Total Duration',   '~6–8 months',                  '~12–15 months',                  '~8–10 months'],
            ['Key Advantage',    'Fastest delivery',             'Builds team ownership',          'Balanced speed + long-term ownership'],
            ['Key Risk',         'Tech team availability',       'Longer timeline; learning curve','Cross-team coordination required'],
        ]
    )

    doc.add_paragraph()
    h2(doc, '6.2  Claude Development Licences')
    doc.add_paragraph(
        'Additional Claude licences are required for the development team during active development. '
        '1 licence is already available. Claude Max: ~$100 USD/user/month. '
        'Licences are only needed during development — not for end users of the live tool.'
    )
    doc_table(doc,
        ['Scenario', 'Additional Licences Needed', 'Active Dev Period', 'Estimated Licence Cost'],
        [
            ['Scenario 1 — Tech Team',      '1–2 licences', '6–8 months',   '~$600 – $1,600 USD'],
            ['Scenario 2 — Analytics Team', '2–3 licences', '12–15 months', '~$2,400 – $4,500 USD'],
            ['Scenario 3 — Hybrid ★',       '2 licences',   '8–10 months',  '~$1,600 – $2,000 USD'],
        ]
    )

    doc.add_paragraph()
    h2(doc, '6.3  Ongoing Operational Costs  (Post Go-Live)')
    doc.add_paragraph('Based on 10 users × 100 questions/month = 1,000 questions/month = 3,000 API calls/month.')
    doc_table(doc,
        ['Cost Item', 'Monthly', 'Annual', 'Notes'],
        [
            ['Claude API',          '~$35 – $50',        '~$420 – $600',   'Only external cost; scales with usage volume'],
            ['Hosting',             '$0',                 '$0',              'Internal BA subdomain'],
            ['Analytics Database',  '$0',                 '$0',              'Existing KpiReports SQL Server'],
            ['Application Database','$0',                 '$0',              'New DB on existing SQL Server infrastructure'],
            ['Maintenance',         'Internal effort',    'Internal effort', 'Owned and maintained by Analytics team'],
            ['TOTAL',               '~$35 – $50/month',  '~$420 – $600/yr', 'Highly cost-effective for an enterprise tool'],
        ],
        highlight_last=True
    )

    # ── 7. RECOMMENDATION ───────────────────────────────────────────────────
    doc.add_paragraph()
    h1(doc, '7.  Recommendation')
    p = doc.add_paragraph()
    r = p.add_run('Scenario 3 — Hybrid Model — is the recommended approach.')
    r.bold = True; r.font.size = DocxPt(12); r.font.color.rgb = D_TEAL

    doc.add_paragraph()
    reasons = [
        ('Domain Ownership',
         'The Analytics team understands the business KPIs, data rules, and what questions matter most. '
         'They are best placed to define what the tool should do — and to maintain it long-term.'),
        ('Claude AI Reduces the Learning Curve',
         'Claude Code significantly lowers the technical barrier. Python, FastAPI, and React '
         'development can be done effectively by analysts with AI assistance — without years of software engineering experience.'),
        ('Bounded Technology Team Involvement',
         'The Tech team\'s role is limited to infrastructure — hosting, database setup, and deployment. '
         'This is a defined, bounded scope that does not create an ongoing dependency.'),
        ('Long-Term Sustainability',
         'A tool expected to be in use for many years should be owned by the team closest to the data. '
         'The hybrid model builds this ownership from day one.'),
        ('Best Value',
         'Fewer licences than Scenario 2, less dependency on Tech bandwidth than Scenario 1, '
         'and faster delivery than the pure analytics team approach.'),
    ]
    for b, d in reasons:
        bullet(doc, b, d)

    # ── 8. NEXT STEPS ───────────────────────────────────────────────────────
    doc.add_paragraph()
    h1(doc, '8.  Next Steps')
    for s in [
        'Review and approve this proposal',
        'Select the implementation scenario',
        'Identify team members — Analytics leads and Technology support contacts',
        'Procure additional Claude development licences for the development team',
        'Set up the application database (user management, access control, query history)',
        'Configure hosting on paneliq.borderlessaccess.com',
        'Complete Phase 1 deployment — POC is ~80% ready, timeline is shorter than a fresh build',
        'Define user access matrix — which users can access which data domains and fields',
    ]:
        p = doc.add_paragraph(style='List Number')
        p.add_run(s).font.size = DocxPt(10.5)

    doc.save(path)
    print(f'  Word document saved: {path}')


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ══════════════════════════════════════════════════════════════════════════════

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def slide_header(slide, prs, title):
    """Dark-blue header bar + teal underline. Returns y offset after header."""
    W = prs.slide_width
    hh = Inches(1.05)

    bar = slide.shapes.add_shape(1, 0, 0, W, hh)
    bar.fill.solid(); bar.fill.fore_color.rgb = P_DARK_BLUE
    bar.line.fill.background()

    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = P_WHITE

    line = slide.shapes.add_shape(1, 0, hh, W, Pt(4))
    line.fill.solid(); line.fill.fore_color.rgb = P_TEAL
    line.line.fill.background()

    return hh + Pt(4)

def txb(slide, text, l, t, w, h, size=11, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color if color else P_GRAY
    return box

def rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        if line_width:
            s.line.width = line_width
    else:
        s.line.fill.background()
    return s

def ppt_table(slide, headers, rows, l, t, w, h,
              first_col_header=False, highlight_col=None, highlight_last_row=False):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), l, t, w, h).table

    for ci, hdr in enumerate(headers):
        c = tbl.cell(0, ci)
        c.fill.solid(); c.fill.fore_color.rgb = P_DARK_BLUE
        r = c.text_frame.paragraphs[0].add_run()
        r.text = hdr; r.font.bold = True
        r.font.color.rgb = P_WHITE; r.font.size = Pt(10)

    for ri, row in enumerate(rows):
        is_last = highlight_last_row and ri == len(rows) - 1
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            if is_last:
                c.fill.solid(); c.fill.fore_color.rgb = P_DARK_BLUE
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.bold = True
                r.font.color.rgb = P_WHITE; r.font.size = Pt(10)
            elif first_col_header and ci == 0:
                c.fill.solid(); c.fill.fore_color.rgb = P_DARK_BLUE
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.bold = True
                r.font.color.rgb = P_WHITE; r.font.size = Pt(10)
            elif highlight_col is not None and ci == highlight_col:
                c.fill.solid(); c.fill.fore_color.rgb = PPTRGB(0xe8, 0xf8, 0xf0)
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.size = Pt(10)
                r.font.color.rgb = P_DARK_BLUE
            elif ri % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = P_LIGHT
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.size = Pt(10); r.font.color.rgb = P_GRAY
            else:
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.size = Pt(10); r.font.color.rgb = P_GRAY
    return tbl


def create_ppt(path):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: TITLE ───────────────────────────────────────────────────────
    s = blank_slide(prs)
    bg = rect(s, 0, 0, prs.slide_width, prs.slide_height, P_DARK_BLUE)
    rect(s, 0, Inches(4.3), prs.slide_width, Pt(5), P_TEAL)

    txb(s, 'PanelIQ', Inches(1.5), Inches(1.3), Inches(10), Inches(1.5),
        size=56, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    txb(s, 'AI-Powered Panel Intelligence Platform',
        Inches(1.5), Inches(2.85), Inches(10), Inches(0.8),
        size=22, color=P_TEAL, align=PP_ALIGN.CENTER)
    txb(s, 'Proposal & Implementation Plan  |  Borderless Access  |  May 2026',
        Inches(1.5), Inches(4.7), Inches(10), Inches(0.6),
        size=13, color=P_MID_GRAY, align=PP_ALIGN.CENTER)

    # ── SLIDE 2: AGENDA ──────────────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Agenda')
    items = [
        ('01', 'Current State & Challenges'),
        ('02', 'Proposed Solution — PanelIQ'),
        ('03', 'Access Control & Security'),
        ('04', 'Technology Stack'),
        ('05', 'Phased Roadmap'),
        ('06', 'Cost & Timeline — 3 Scenarios'),
        ('07', 'Recommendation'),
        ('08', 'Next Steps'),
    ]
    for i, (num, label) in enumerate(items):
        col = i // 4
        row = i % 4
        lft = Inches(0.6) + col * Inches(6.4)
        top = ct + Inches(0.3) + row * Inches(1.4)
        box = rect(s, lft, top, Inches(5.9), Inches(1.2), P_LIGHT, P_TEAL, Pt(1))
        txb(s, f'{num}  {label}', lft + Inches(0.2), top + Inches(0.3),
            Inches(5.5), Inches(0.7), size=14, bold=True, color=P_DARK_BLUE)

    # ── SLIDE 3: CURRENT STATE ───────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Current State & Challenges')

    txb(s, 'How Insights Are Accessed Today', Inches(0.5), ct + Inches(0.15),
        Inches(5.8), Inches(0.4), size=13, bold=True, color=P_TEAL)

    today = [
        'Tableau dashboards — pre-defined, static performance views',
        'SSRS reports — structured operational reporting',
        'Ad-hoc requests to Technology, Finance, or Analytics teams',
        '1–3 business days wait for any dynamic question',
        'Dynamic insights available only at monthly/quarterly reviews',
    ]
    for i, item in enumerate(today):
        txb(s, f'• {item}',
            Inches(0.5), ct + Inches(0.65) + i * Inches(1.0),
            Inches(5.8), Inches(0.85), size=11, color=P_GRAY)

    rect(s, Inches(6.5), ct + Inches(0.1), Pt(3), Inches(6.1), P_TEAL)

    txb(s, 'Key Challenges', Inches(6.7), ct + Inches(0.15),
        Inches(6.2), Inches(0.4), size=13, bold=True, color=P_TEAL)

    challenges = [
        ('Delayed Decisions',   '1–3 days for any dynamic question'),
        ('Team Dependency',     'Every request creates workload for support teams'),
        ('Low Frequency',       'Insights available monthly or quarterly only'),
        ('Inconsistency Risk',  'Different teams, different metric interpretations'),
        ('No Self-Service',     'Non-technical users locked out of dynamic data'),
        ('Scalability Gap',     'Request volume will grow with the business'),
    ]
    for i, (t, d) in enumerate(challenges):
        bg_c = P_LIGHT if i % 2 == 0 else P_WHITE
        bx = rect(s, Inches(6.7), ct + Inches(0.65) + i * Inches(0.95),
                  Inches(6.2), Inches(0.85), bg_c)
        box = slide_header  # reuse helper name in scope — don't need it
        tb = s.shapes.add_textbox(
            Inches(6.9), ct + Inches(0.7) + i * Inches(0.95),
            Inches(5.9), Inches(0.75)
        )
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = t + '  '
        r1.font.bold = True; r1.font.size = Pt(10.5); r1.font.color.rgb = P_DARK_BLUE
        r2 = p.add_run(); r2.text = d
        r2.font.size = Pt(10.5); r2.font.color.rgb = P_GRAY

    # ── SLIDE 4: SOLUTION ────────────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Proposed Solution — PanelIQ')

    txb(s,
        'A secure web-based AI assistant at paneliq.borderlessaccess.com — '
        'ask in plain English, get answers in seconds.',
        Inches(0.5), ct + Inches(0.1), Inches(12.3), Inches(0.45),
        size=11.5, italic=True, color=P_GRAY)

    steps = [
        ('1', 'User Asks',         '"Show me completes by\nmarket for Q1 2026"'),
        ('2', 'Claude Generates SQL',  'Translates natural language\ninto precise, secure SQL'),
        ('3', 'Query Executes',    'Runs securely against\nKpiReports database'),
        ('4', 'Instant Results',   'Chart + table + analyst\nsummary in seconds'),
    ]
    for i, (num, title, desc) in enumerate(steps):
        lft = Inches(0.4) + i * Inches(3.25)
        tp  = ct + Inches(0.65)
        rect(s, lft, tp, Inches(3.05), Inches(4.2), P_LIGHT, P_TEAL, Pt(1))

        circ = s.shapes.add_shape(9, lft + Inches(1.0), tp + Inches(0.25),
                                   Inches(1.0), Inches(1.0))
        circ.fill.solid(); circ.fill.fore_color.rgb = P_TEAL
        circ.line.fill.background()
        ctf = circ.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(20); cr.font.bold = True; cr.font.color.rgb = P_WHITE

        txb(s, title, lft + Inches(0.1), tp + Inches(1.4), Inches(2.85), Inches(0.55),
            size=12, bold=True, color=P_DARK_BLUE, align=PP_ALIGN.CENTER)
        txb(s, desc, lft + Inches(0.1), tp + Inches(2.05), Inches(2.85), Inches(1.9),
            size=10.5, color=P_GRAY, align=PP_ALIGN.CENTER)

        if i < 3:
            txb(s, '→', Inches(3.4) + i * Inches(3.25), ct + Inches(2.2),
                Inches(0.3), Inches(0.5), size=22, bold=True, color=P_TEAL)

    # POC Banner
    rect(s, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.7), P_GREEN)
    txb(s, '✓  Proof of Concept validated  —  Phase 1 is approximately 80% complete',
        Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
        size=12, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    # ── SLIDE 5: ACCESS CONTROL ──────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Access Control & Data Security')

    panels = [
        ('Role-Based Access',
         'Users see only data relevant to their role and BU. A market analyst sees their market; '
         'a BU director sees their BU.',
         P_TEAL),
        ('Field-Level Restrictions',
         'Sensitive fields — revenue, vendor cost, margins — visible only to authorised users. '
         'Finance sees financials; operations see operations.',
         P_DARK_BLUE),
        ('Row-Level Scoping',
         'A BU1 director sees only BU1 client revenue. BU2 data is never included in the query — '
         'not filtered after retrieval.',
         P_TEAL),
        ('AI-Enforced Rules',
         'Access constraints are injected into every SQL query at the AI layer. '
         'Users cannot bypass restrictions through creative questioning.',
         P_DARK_BLUE),
    ]
    for i, (title, desc, color) in enumerate(panels):
        col = i % 2; row = i // 2
        lft = Inches(0.4) + col * Inches(6.45)
        tp  = ct + Inches(0.25) + row * Inches(2.8)
        bx = rect(s, lft, tp, Inches(6.1), Inches(2.6), P_LIGHT, color, Pt(2))
        txb(s, title, lft + Inches(0.2), tp + Inches(0.2),
            Inches(5.7), Inches(0.5), size=13, bold=True, color=color)
        txb(s, desc, lft + Inches(0.2), tp + Inches(0.8),
            Inches(5.7), Inches(1.6), size=11, color=P_GRAY, wrap=True)

    # ── SLIDE 6: TECH STACK ──────────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Technology Stack')
    txb(s,
        'Enterprise-grade technologies selected for long-term scalability, developer availability, and multi-year support',
        Inches(0.5), ct + Inches(0.05), Inches(12.3), Inches(0.4),
        size=11, italic=True, color=P_GRAY)
    ppt_table(s,
        ['Component', 'Technology', 'Purpose', 'Scalability'],
        [
            ['AI Engine',      'Anthropic Claude API',     'NL-to-SQL, summaries, suggestions',      'API-based; no infra; model improves over time'],
            ['Backend',        'Python + FastAPI',          'Query engine, logic, access control',     'Production-grade; high concurrency'],
            ['Frontend',       'React + Vite',              'Chat UI, charts, data tables',            'Industry standard; large ecosystem'],
            ['Analytics DB',   'SQL Server (KpiReports)',   'Source of truth for panel data',          'Enterprise-grade; existing infrastructure'],
            ['App DB',         'SQL Server (new internal)', 'Users, history, access rules',            'Existing SQL Server — no added licence cost'],
            ['Hosting',        'Internal subdomain',        'paneliq.borderlessaccess.com',            'No cost; scales to 500+ users'],
            ['Version Control','GitHub',                    'Code management, team collaboration',     'Industry standard; private repository'],
            ['Dev Assistant',  'Claude Code (Anthropic)',   'AI-assisted development',                 'Greatly reduces learning curve for analysts'],
        ],
        Inches(0.4), ct + Inches(0.5), Inches(12.5), Inches(5.7),
        first_col_header=True
    )

    # ── SLIDE 7: ROADMAP ─────────────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Phased Roadmap')

    phases = [
        ('Phase 1', 'Foundation',          '~80% Complete',        P_GREEN,
         ['Hosted deployment', 'User auth & login', 'Role + row-level access control',
          'Sample analytics (markets, vendors, IR%, rates)', 'Query history & suggestions', 'App DB setup']),
        ('Phase 2', 'Vendor & Quality',    'Pending schema',       P_TEAL,
         ['Vendor cost analysis', 'Sample quality metrics', 'Vendor scoring & comparison', 'New SQL tables integration']),
        ('Phase 3', 'Revenue & Commercial','Pending Finance data',  P_DARK_BLUE,
         ['Revenue by market/client/BU', 'Research cost & margins', 'Financial access control', 'Finance data integration']),
        ('Phase 4', 'Advanced Intelligence','Future',              P_PURPLE,
         ['Trend forecasting', 'Anomaly alerts', 'Executive summaries', 'Tableau/SSRS integration']),
    ]
    for i, (ph, title, status, color, items) in enumerate(phases):
        lft = Inches(0.3) + i * Inches(3.25)
        tp  = ct + Inches(0.2)

        rect(s, lft, tp, Inches(3.1), Inches(0.75), color)
        txb(s, ph,    lft + Inches(0.1), tp + Inches(0.02), Inches(2.9), Inches(0.3),
            size=10, bold=True, color=P_WHITE)
        txb(s, title, lft + Inches(0.1), tp + Inches(0.32), Inches(2.9), Inches(0.4),
            size=13, bold=True, color=P_WHITE)

        rect(s, lft, tp + Inches(0.75), Inches(3.1), Inches(0.42), P_LIGHT, color, Pt(1))
        txb(s, status, lft + Inches(0.1), tp + Inches(0.78),
            Inches(2.9), Inches(0.35), size=10, italic=True, color=color, align=PP_ALIGN.CENTER)

        rect(s, lft, tp + Inches(1.17), Inches(3.1), Inches(5.0),
             P_WHITE, PPTRGB(0xd0, 0xde, 0xe8), Pt(1))
        for j, item in enumerate(items):
            txb(s, f'• {item}',
                lft + Inches(0.15), tp + Inches(1.28) + j * Inches(0.72),
                Inches(2.8), Inches(0.65), size=10.5, color=P_GRAY)

    # ── SLIDE 8: SCENARIOS ───────────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Implementation Scenarios — Cost & Timeline')

    ppt_table(s,
        ['', 'Scenario 1\nTech Team', 'Scenario 2\nAnalytics Team', 'Scenario 3  ★\nHybrid (Recommended)'],
        [
            ['Team',              '1–2 Tech developers',       '2–3 Analysts + Claude AI',   'Analytics leads + Tech for infra'],
            ['Phase 1',           '4–6 weeks',                 '10–14 weeks',                 '5–7 weeks'],
            ['Phase 2',           '6–8 weeks',                 '10–14 weeks',                 '7–10 weeks'],
            ['Phase 3',           '6–8 weeks',                 '10–14 weeks',                 '7–10 weeks'],
            ['Phase 4',           '8–10 weeks',                '14–18 weeks',                 '10–12 weeks'],
            ['Total Duration',    '~6–8 months',               '~12–15 months',               '~8–10 months'],
            ['Claude Licences',   '2 × ~$100/mo',              '3 × ~$100/mo',                '2 × ~$100/mo'],
            ['Licence Cost (dev)','~$1,200–$1,600',            '~$2,400–$4,500',              '~$1,600–$2,000'],
            ['Key Advantage',     'Fastest delivery',          'Full team ownership',         'Balanced speed + ownership'],
            ['Key Risk',          'Tech team availability',    'Longer timeline',             'Cross-team coordination'],
        ],
        Inches(0.4), ct + Inches(0.15), Inches(12.5), Inches(6.1),
        first_col_header=True, highlight_col=3
    )

    # ── SLIDE 9: OPERATIONAL COSTS ───────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Ongoing Operational Costs  (Post Go-Live)')

    txb(s,
        'Based on 10 users × 100 questions/month = 1,000 questions/month = 3,000 API calls/month',
        Inches(0.5), ct + Inches(0.05), Inches(12.3), Inches(0.4),
        size=11, italic=True, color=P_GRAY)

    ppt_table(s,
        ['Cost Item', 'Monthly', 'Annual', 'Notes'],
        [
            ['Claude API',           '~$35 – $50',        '~$420 – $600',    'Only external cost — scales with usage'],
            ['Hosting',              '$0',                 '$0',               'Internal BA subdomain'],
            ['Analytics Database',   '$0',                 '$0',               'Existing KpiReports SQL Server'],
            ['Application Database', '$0',                 '$0',               'New DB on existing SQL Server infra'],
            ['Maintenance',          'Internal',           'Internal',         'Owned by Analytics team'],
            ['TOTAL',                '~$35–$50/month',     '~$420–$600/year',  'Highly cost-effective for an enterprise tool'],
        ],
        Inches(0.8), ct + Inches(0.55), Inches(11.6), Inches(4.3),
        highlight_last_row=True
    )

    rect(s, Inches(0.8), Inches(6.05), Inches(11.6), Inches(0.82),
         PPTRGB(0xe8, 0xf8, 0xf0), P_GREEN, Pt(1))
    txb(s,
        'The only recurring external cost is the Anthropic API  —  under $50/month for a 10-user deployment',
        Inches(0.9), Inches(6.1), Inches(11.4), Inches(0.7),
        size=12, bold=True, color=PPTRGB(0x00, 0x6b, 0x45), align=PP_ALIGN.CENTER)

    # ── SLIDE 10: RECOMMENDATION ─────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Recommendation — Scenario 3: Hybrid Model')

    txb(s,
        'Analytics team leads  ·  Technology team provides infrastructure support  ·  Claude AI accelerates delivery',
        Inches(0.5), ct + Inches(0.1), Inches(12.3), Inches(0.4),
        size=12, italic=True, color=P_TEAL)

    reasons = [
        ('Domain Ownership',
         'Analytics team understands KPIs and data rules. Best placed to build and maintain long-term.',
         P_TEAL),
        ('Claude AI Lowers the Bar',
         'Python, FastAPI, React — all achievable by analysts with Claude Code assistance.',
         P_DARK_BLUE),
        ('Bounded Tech Involvement',
         'Tech team role: hosting, DB setup, deployment. Defined scope. No ongoing dependency.',
         P_TEAL),
        ('Long-Term Sustainability',
         'Tool ownership stays with the team closest to the data from day one.',
         P_DARK_BLUE),
        ('Best Value',
         'Fewer licences than Scenario 2. Less Tech dependency than Scenario 1. Faster delivery.',
         P_GREEN),
    ]
    for i, (title, desc, color) in enumerate(reasons):
        col = i % 2 if i < 4 else 0
        row = i // 2
        if i == 4:
            lft = Inches(3.7)
        else:
            lft = Inches(0.4) + col * Inches(6.45)
        tp = ct + Inches(0.65) + row * Inches(1.85)
        w  = Inches(12.5) if i == 4 else Inches(6.1)

        bx = rect(s, lft, tp, w, Inches(1.7), P_LIGHT, color, Pt(1))
        txb(s, title, lft + Inches(0.2), tp + Inches(0.15), w - Inches(0.4), Inches(0.45),
            size=12, bold=True, color=color)
        txb(s, desc,  lft + Inches(0.2), tp + Inches(0.65), w - Inches(0.4), Inches(0.95),
            size=10.5, color=P_GRAY, wrap=True)

    # ── SLIDE 11: NEXT STEPS ─────────────────────────────────────────────────
    s = blank_slide(prs)
    ct = slide_header(s, prs, 'Next Steps')

    steps = [
        ('1', 'Approve Proposal',      'Review, align on scope and scenario',                                P_TEAL),
        ('2', 'Select Scenario',       'Confirm team structure — hybrid model recommended',                   P_DARK_BLUE),
        ('3', 'Identify Team',         'Name Analytics leads and Technology support contacts',               P_TEAL),
        ('4', 'Procure Licences',      'Additional Claude development licences for the team',                P_DARK_BLUE),
        ('5', 'Setup App Database',    'User management, access control rules, query history storage',       P_TEAL),
        ('6', 'Configure Hosting',     'Set up paneliq.borderlessaccess.com subdomain',                      P_DARK_BLUE),
        ('7', 'Phase 1 Deployment',    'POC is ~80% ready — faster than a fresh build',                      P_TEAL),
        ('8', 'Define Access Matrix',  'Map users to data domains, roles, and field permissions',            P_DARK_BLUE),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        col = i % 2; row = i // 2
        lft = Inches(0.4) + col * Inches(6.45)
        tp  = ct + Inches(0.25) + row * Inches(1.5)

        circ = s.shapes.add_shape(9, lft, tp + Inches(0.3),
                                   Inches(0.65), Inches(0.65))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        ctf = circ.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(14); cr.font.bold = True; cr.font.color.rgb = P_WHITE

        tb = s.shapes.add_textbox(lft + Inches(0.8), tp + Inches(0.15),
                                   Inches(5.4), Inches(1.2))
        tf = tb.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = title
        r1.font.bold = True; r1.font.size = Pt(12); r1.font.color.rgb = color
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = desc
        r2.font.size = Pt(10.5); r2.font.color.rgb = P_GRAY

    # ── SLIDE 12: THANK YOU ──────────────────────────────────────────────────
    s = blank_slide(prs)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, P_DARK_BLUE)
    rect(s, 0, Inches(3.5), prs.slide_width, Pt(5), P_TEAL)

    txb(s, 'Thank You',
        Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
        size=52, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    txb(s, 'Questions & Discussion',
        Inches(1.5), Inches(3.0), Inches(10), Inches(0.7),
        size=20, color=P_TEAL, align=PP_ALIGN.CENTER)
    txb(s, 'paneliq.borderlessaccess.com  |  Borderless Access  |  May 2026',
        Inches(1.5), Inches(4.3), Inches(10), Inches(0.6),
        size=13, color=P_MID_GRAY, align=PP_ALIGN.CENTER)

    prs.save(path)
    print(f'  PowerPoint saved: {path}')


# ── RUN ───────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    create_word(os.path.join(BASE, 'PanelIQ_Proposal_May2026.docx'))
    create_ppt(os.path.join(BASE,  'PanelIQ_Proposal_May2026.pptx'))
    print('Done.')
