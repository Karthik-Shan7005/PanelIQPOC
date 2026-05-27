"""
generate_ba_ppts.py
Regenerates both PanelIQ PPT files using BA Template 2026.pptx as the base.
Master slide carries logo, footer, and bottom bars automatically.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RGB
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from io import BytesIO
import os

TEMPLATE = r'C:\Users\KarthikShanmugam\ClaudePOC\BA Template 2026.pptx'
BASE     = r'C:\Users\KarthikShanmugam\ClaudePOC\paneliq'

# ── BA brand colours ──────────────────────────────────────────────────────────
BA_BLUE   = RGB(0x03, 0x6d, 0xac)   # primary blue
BA_NAVY   = RGB(0x0e, 0x28, 0x41)   # dark navy
BA_TEAL   = RGB(0x00, 0xae, 0xae)   # teal
BA_ORANGE = RGB(0xe9, 0x71, 0x32)   # orange accent
BA_LBLUE  = RGB(0x0f, 0x9e, 0xd5)   # light blue accent
BA_GREEN  = RGB(0x00, 0xa8, 0x6b)   # green
BA_AMBER  = RGB(0xf2, 0xb9, 0x0c)   # amber
WHITE     = RGB(0xff, 0xff, 0xff)
GRAY      = RGB(0x2d, 0x37, 0x48)
LIGHT     = RGB(0xf0, 0xf7, 0xfb)   # very light blue
MID_GRAY  = RGB(0xa0, 0xb8, 0xd0)
PALE_GREEN= RGB(0xe8, 0xf8, 0xf0)
PALE_BLUE = RGB(0xe8, 0xf4, 0xfb)
RED_SOFT  = RGB(0xc0, 0x39, 0x2b)


# ── Helpers ───────────────────────────────────────────────────────────────────

def load_template():
    """Open BA template and properly remove all existing slides."""
    prs = Presentation(TEMPLATE)
    prs_xml = prs.part._element
    sldIdLst = prs_xml.find(qn('p:sldIdLst'))

    # Collect rIds then drop both the sldIdLst entry AND the OPC relationship,
    # so the old slide parts are not saved into the output zip.
    rIds = [sldId.get(qn('r:id')) for sldId in list(sldIdLst)]
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)
    for rId in rIds:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[47])  # 'Blank'

def _logo_bytes(prs):
    """Extract BA logo PNG from slide master (Picture 10)."""
    for sh in prs.slide_master.shapes:
        if sh.shape_type == 13:           # MSO_SHAPE_TYPE.PICTURE = 13
            try:
                return sh.image.blob
            except Exception:
                pass
    return None

def add_logo(slide, prs):
    """Re-pin BA logo at its master position so it sits above colored bars."""
    blob = _logo_bytes(prs)
    if blob:
        slide.shapes.add_picture(
            BytesIO(blob),
            Inches(11.767), Inches(0.099),
            Inches(1.496), Inches(0.473)
        )

def rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw:
            s.line.width = lw
    else:
        s.line.fill.background()
    return s

def oval(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txb(slide, text, l, t, w, h,
        size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    s = slide.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = 'Arial'
    r.font.color.rgb = color if color else GRAY
    return s

def two_run(slide, l, t, w, h, t1, s1, b1, c1, t2, s2, b2, c2,
            align=PP_ALIGN.LEFT):
    s = slide.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r1 = p.add_run()
    r1.text = t1
    r1.font.size = Pt(s1); r1.font.bold = b1
    r1.font.color.rgb = c1; r1.font.name = 'Arial'
    r2 = p.add_run()
    r2.text = t2
    r2.font.size = Pt(s2); r2.font.bold = b2
    r2.font.color.rgb = c2; r2.font.name = 'Arial'

def slide_header(slide, prs, title):
    """BA blue header bar with teal underline + logo re-pinned on top."""
    rect(slide, 0, 0, prs.slide_width, Inches(1.0), BA_BLUE)
    rect(slide, 0, Inches(1.0), prs.slide_width, Pt(3), BA_TEAL)
    txb(slide, title,
        Inches(0.4), Inches(0.15), Inches(11.1), Inches(0.72),
        size=22, bold=True, color=WHITE)
    add_logo(slide, prs)
    return Inches(1.07)

def stat_block(slide, l, t, w, h, big, label, bg, big_c, lbl_c):
    rect(slide, l, t, w, h, bg, BA_BLUE, Pt(1))
    txb(slide, big,   l, t + Inches(0.2),  w, Inches(0.85),
        size=40, bold=True, color=big_c, align=PP_ALIGN.CENTER)
    txb(slide, label, l, t + Inches(1.0),  w, Inches(0.95),
        size=11, color=lbl_c, align=PP_ALIGN.CENTER, wrap=True)

def ppt_table(slide, headers, rows, l, t, w, h,
              first_col_header=False, highlight_col=None, highlight_last_row=False):
    tbl = slide.shapes.add_table(len(rows)+1, len(headers), l, t, w, h).table
    for ci, hdr in enumerate(headers):
        c = tbl.cell(0, ci)
        c.fill.solid(); c.fill.fore_color.rgb = BA_BLUE
        r = c.text_frame.paragraphs[0].add_run()
        r.text = hdr; r.font.bold = True; r.font.name = 'Arial'
        r.font.color.rgb = WHITE; r.font.size = Pt(10)

    for ri, row in enumerate(rows):
        is_last = highlight_last_row and ri == len(rows)-1
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci)
            if is_last:
                c.fill.solid(); c.fill.fore_color.rgb = BA_NAVY
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.bold = True; r.font.name = 'Arial'
                r.font.color.rgb = WHITE; r.font.size = Pt(10)
            elif first_col_header and ci == 0:
                c.fill.solid(); c.fill.fore_color.rgb = BA_NAVY
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.bold = True; r.font.name = 'Arial'
                r.font.color.rgb = WHITE; r.font.size = Pt(10)
            elif highlight_col is not None and ci == highlight_col:
                c.fill.solid(); c.fill.fore_color.rgb = PALE_GREEN
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.size = Pt(10); r.font.name = 'Arial'
                r.font.color.rgb = BA_NAVY
            elif ri % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.size = Pt(10); r.font.name = 'Arial'; r.font.color.rgb = GRAY
            else:
                r = c.text_frame.paragraphs[0].add_run()
                r.text = val; r.font.size = Pt(10); r.font.name = 'Arial'; r.font.color.rgb = GRAY
    return tbl


# ══════════════════════════════════════════════════════════════════════════════
#  DETAILED PROPOSAL DECK  (12 slides)
# ══════════════════════════════════════════════════════════════════════════════

def detailed_ppt(path):
    prs = load_template()
    W = prs.slide_width

    # ── S1: TITLE ─────────────────────────────────────────────────────────────
    s = blank(prs)
    # Navy panel left 65%
    rect(s, 0, 0, Inches(8.7), prs.slide_height, BA_NAVY)
    # Teal accent bar
    rect(s, 0, Inches(4.2), Inches(8.7), Pt(4), BA_TEAL)
    # White right panel (so master logo is visible)
    rect(s, Inches(8.7), 0, Inches(4.63), prs.slide_height, WHITE)

    txb(s, 'PanelIQ',
        Inches(0.5), Inches(1.3), Inches(7.9), Inches(1.4),
        size=52, bold=True, color=WHITE)
    txb(s, 'AI-Powered Panel Intelligence Platform',
        Inches(0.5), Inches(2.85), Inches(7.9), Inches(0.7),
        size=20, color=BA_TEAL)
    txb(s, 'Proposal & Implementation Plan  |  Borderless Access  |  May 2026',
        Inches(0.5), Inches(4.5), Inches(7.9), Inches(0.55),
        size=13, color=MID_GRAY)

    # BA blue vertical accent on left edge
    rect(s, 0, 0, Inches(0.18), prs.slide_height, BA_BLUE)

    # Right panel: tagline
    txb(s, 'From waiting days\nto answers in seconds.',
        Inches(9.0), Inches(2.8), Inches(4.0), Inches(1.4),
        size=15, bold=True, color=BA_BLUE, wrap=True)

    # ── S2: AGENDA ────────────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Agenda')
    items = [
        ('01', 'Current State & Challenges'),
        ('02', 'Proposed Solution — PanelIQ'),
        ('03', 'Access Control & Security'),
        ('04', 'Technology Stack'),
        ('05', 'Phased Roadmap'),
        ('06', 'Cost & Timeline — 3 Scenarios'),
        ('07', 'Recommendation'),
        ('08', 'Next Steps'),
    ]
    for i, (num, label) in enumerate(items):
        col = i // 4; row = i % 4
        lft = Inches(0.5) + col * Inches(6.4)
        tp  = ct + Inches(0.3) + row * Inches(1.35)
        rect(s, lft, tp, Inches(6.0), Inches(1.2), LIGHT, BA_BLUE, Pt(1))
        txb(s, f'{num}  {label}', lft + Inches(0.2), tp + Inches(0.3),
            Inches(5.6), Inches(0.65), size=14, bold=True, color=BA_BLUE)

    # ── S3: CURRENT STATE ─────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Current State & Challenges')

    txb(s, 'How Insights Are Accessed Today',
        Inches(0.4), ct + Inches(0.15), Inches(5.9), Inches(0.4),
        size=13, bold=True, color=BA_BLUE)
    today = [
        'Tableau dashboards — pre-defined, static views only',
        'SSRS reports — structured operational reporting',
        'Ad-hoc requests to Technology, Finance, or Analytics teams',
        '1–3 business days wait for any dynamic question',
        'Dynamic insights only at monthly/quarterly reviews',
    ]
    for i, item in enumerate(today):
        txb(s, f'• {item}', Inches(0.4),
            ct + Inches(0.65) + i * Inches(1.0), Inches(5.9), Inches(0.85),
            size=11, color=GRAY)

    rect(s, Inches(6.5), ct + Inches(0.1), Pt(3), Inches(5.6), BA_TEAL)

    txb(s, 'Key Challenges',
        Inches(6.7), ct + Inches(0.15), Inches(6.2), Inches(0.4),
        size=13, bold=True, color=BA_BLUE)
    challenges = [
        ('Delayed Decisions',   '1–3 days for any dynamic question'),
        ('Team Dependency',     'Every request creates workload for support teams'),
        ('Low Frequency',       'Insights available monthly or quarterly only'),
        ('Inconsistency Risk',  'Different teams may interpret the same metric differently'),
        ('No Self-Service',     'Non-technical users locked out of dynamic data'),
        ('Scalability Gap',     'Request volume grows with the business'),
    ]
    for i, (title, desc) in enumerate(challenges):
        bg = LIGHT if i % 2 == 0 else WHITE
        rect(s, Inches(6.7), ct + Inches(0.65) + i * Inches(0.92),
             Inches(6.2), Inches(0.82), bg)
        tb = s.shapes.add_textbox(
            Inches(6.9), ct + Inches(0.68) + i * Inches(0.92),
            Inches(5.9), Inches(0.72))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = title + '  '
        r1.font.bold = True; r1.font.size = Pt(10.5)
        r1.font.color.rgb = BA_NAVY; r1.font.name = 'Arial'
        r2 = p.add_run(); r2.text = desc
        r2.font.size = Pt(10.5); r2.font.color.rgb = GRAY; r2.font.name = 'Arial'

    # ── S4: SOLUTION ──────────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Proposed Solution — PanelIQ')

    txb(s, 'A secure web-based AI assistant at paneliq.borderlessaccess.com — '
           'ask in plain English, get answers in seconds.',
        Inches(0.4), ct + Inches(0.1), Inches(12.5), Inches(0.45),
        size=11.5, italic=True, color=GRAY)

    steps = [
        ('1', 'User Asks',          '"Show me completes by\nmarket for Q1 2026"',    BA_BLUE),
        ('2', 'Claude Generates SQL','Translates natural language\ninto precise SQL', BA_NAVY),
        ('3', 'Query Executes',      'Runs securely against\nKpiReports database',    BA_TEAL),
        ('4', 'Instant Results',     'Chart + table + analyst\nsummary in seconds',   BA_GREEN),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        lft = Inches(0.35) + i * Inches(3.25)
        tp  = ct + Inches(0.65)
        rect(s, lft, tp, Inches(3.05), Inches(4.0), LIGHT, color, Pt(1))
        c = oval(s, lft + Inches(1.0), tp + Inches(0.22), Inches(1.05), Inches(1.05), color)
        ctf = c.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(22); cr.font.bold = True
        cr.font.color.rgb = WHITE; cr.font.name = 'Arial'
        txb(s, title, lft + Inches(0.1), tp + Inches(1.4),
            Inches(2.85), Inches(0.5), size=12, bold=True, color=color,
            align=PP_ALIGN.CENTER)
        txb(s, desc, lft + Inches(0.1), tp + Inches(2.0),
            Inches(2.85), Inches(1.8), size=10.5, color=GRAY,
            align=PP_ALIGN.CENTER, wrap=True)
        if i < 3:
            txb(s, '→', Inches(3.35) + i * Inches(3.25), ct + Inches(2.35),
                Inches(0.3), Inches(0.5), size=22, bold=True, color=BA_TEAL)

    rect(s, Inches(0.35), Inches(6.1), Inches(12.6), Inches(0.72), BA_GREEN)
    txb(s, '✓  Proof of Concept validated  —  Phase 1 is approximately 80% complete',
        Inches(0.45), Inches(6.15), Inches(12.4), Inches(0.6),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── S5: ACCESS CONTROL ────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Access Control & Data Security')

    panels = [
        ('Role-Based Access',
         'Users see only data relevant to their role and BU. '
         'A market analyst sees their market; a BU director sees their BU.',
         BA_BLUE),
        ('Field-Level Restrictions',
         'Sensitive fields — revenue, vendor cost, margins — visible only to authorised users. '
         'Finance sees financials; operations see operations.',
         BA_NAVY),
        ('Row-Level Scoping',
         'A BU1 director sees only BU1 client revenue. '
         'BU2 data is never included in the query — not filtered after retrieval.',
         BA_BLUE),
        ('AI-Enforced Rules',
         'Access constraints are injected into every SQL query at the AI layer. '
         'Users cannot bypass restrictions through conversation.',
         BA_NAVY),
    ]
    for i, (title, desc, color) in enumerate(panels):
        col = i % 2; row = i // 2
        lft = Inches(0.4) + col * Inches(6.45)
        tp  = ct + Inches(0.25) + row * Inches(2.7)
        rect(s, lft, tp, Inches(6.1), Inches(2.55), LIGHT, color, Pt(2))
        txb(s, title, lft + Inches(0.2), tp + Inches(0.2),
            Inches(5.7), Inches(0.5), size=13, bold=True, color=color)
        txb(s, desc, lft + Inches(0.2), tp + Inches(0.8),
            Inches(5.7), Inches(1.55), size=11, color=GRAY, wrap=True)

    # ── S6: TECH STACK ────────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Technology Stack')
    txb(s, 'Enterprise-grade technologies — selected for long-term scalability, '
           'developer availability, and multi-year support',
        Inches(0.4), ct + Inches(0.05), Inches(12.5), Inches(0.4),
        size=11, italic=True, color=GRAY)
    ppt_table(s,
        ['Component', 'Technology', 'Purpose', 'Scalability'],
        [
            ['AI Engine',      'Anthropic Claude API',     'NL-to-SQL, summaries, suggestions',      'API-based; model improves continuously'],
            ['Backend',        'Python + FastAPI',          'Query engine, logic, access control',     'Production-grade; high concurrency'],
            ['Frontend',       'React + Vite',              'Chat UI, charts, data tables',            'Industry standard; large ecosystem'],
            ['Analytics DB',   'SQL Server (KpiReports)',   'Source of truth for all panel data',      'Enterprise-grade; existing infrastructure'],
            ['App DB',         'SQL Server (new internal)', 'Users, history, access rules',            'Existing SQL Server — no added licence'],
            ['Hosting',        'Internal subdomain',        'paneliq.borderlessaccess.com',            'No cost; scales to 500+ users'],
            ['Version Control','GitHub',                    'Code management, collaboration',          'Industry standard; private repository'],
            ['Dev Assistant',  'Claude Code (Anthropic)',   'AI-assisted development',                 'Reduces learning curve for analysts'],
        ],
        Inches(0.35), ct + Inches(0.5), Inches(12.6), Inches(5.65),
        first_col_header=True)

    # ── S7: ROADMAP ───────────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Phased Roadmap')

    phases = [
        ('Phase 1', 'Foundation',           '~80% Complete',       BA_GREEN,
         ['Hosted deployment', 'User auth & login',
          'Role + row-level access control', 'Sample analytics (markets, vendors, IR%)',
          'Query history & suggestions', 'App DB setup']),
        ('Phase 2', 'Vendor & Quality',     'Pending schema',      BA_TEAL,
         ['Vendor cost analysis', 'Sample quality metrics',
          'Vendor scoring', 'New SQL tables integration']),
        ('Phase 3', 'Revenue & Commercial', 'Pending Finance data', BA_BLUE,
         ['Revenue by market/client/BU', 'Research cost & margins',
          'Financial access control', 'Finance data integration']),
        ('Phase 4', 'Advanced Intelligence','Future',              BA_NAVY,
         ['Trend forecasting', 'Anomaly alerts',
          'Executive summaries', 'Tableau/SSRS integration']),
    ]
    for i, (ph, title, status, color, items) in enumerate(phases):
        lft = Inches(0.3) + i * Inches(3.25)
        tp  = ct + Inches(0.2)
        rect(s, lft, tp, Inches(3.1), Inches(0.72), color)
        txb(s, ph,    lft + Inches(0.1), tp + Inches(0.02),
            Inches(2.9), Inches(0.28), size=10, bold=True, color=WHITE)
        txb(s, title, lft + Inches(0.1), tp + Inches(0.3),
            Inches(2.9), Inches(0.38), size=13, bold=True, color=WHITE)
        rect(s, lft, tp + Inches(0.72), Inches(3.1), Inches(0.42), LIGHT, color, Pt(1))
        txb(s, status, lft + Inches(0.1), tp + Inches(0.76),
            Inches(2.9), Inches(0.35), size=10, italic=True, color=color,
            align=PP_ALIGN.CENTER)
        rect(s, lft, tp + Inches(1.14), Inches(3.1), Inches(5.05),
             WHITE, RGB(0xd0, 0xde, 0xe8), Pt(1))
        for j, item in enumerate(items):
            txb(s, f'• {item}',
                lft + Inches(0.15), tp + Inches(1.25) + j * Inches(0.72),
                Inches(2.8), Inches(0.65), size=10.5, color=GRAY)

    # ── S8: SCENARIOS ─────────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Implementation Scenarios — Cost & Timeline')
    ppt_table(s,
        ['', 'Scenario 1\nTech Team', 'Scenario 2\nAnalytics Team',
         'Scenario 3  ★\nHybrid (Recommended)'],
        [
            ['Team',              '1–2 Tech developers',       '2–3 Analysts + Claude AI',   'Analytics leads + Tech for infra'],
            ['Phase 1',           '4–6 weeks',                 '10–14 weeks',                '5–7 weeks'],
            ['Phase 2',           '6–8 weeks',                 '10–14 weeks',                '7–10 weeks'],
            ['Phase 3',           '6–8 weeks',                 '10–14 weeks',                '7–10 weeks'],
            ['Phase 4',           '8–10 weeks',                '14–18 weeks',                '10–12 weeks'],
            ['Total Duration',    '~6–8 months',               '~12–15 months',              '~8–10 months'],
            ['Claude Licences',   '2 × ~$100/mo',              '3 × ~$100/mo',               '2 × ~$100/mo'],
            ['Licence Cost (dev)','~$1,200–$1,600',            '~$2,400–$4,500',             '~$1,600–$2,000'],
            ['Key Advantage',     'Fastest delivery',          'Full team ownership',        'Balanced speed + ownership'],
            ['Key Risk',          'Tech team availability',    'Longer timeline',            'Cross-team coordination'],
        ],
        Inches(0.35), ct + Inches(0.15), Inches(12.6), Inches(6.1),
        first_col_header=True, highlight_col=3)

    # ── S9: OPERATIONAL COSTS ─────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Ongoing Operational Costs  (Post Go-Live)')
    txb(s, '10 users × 100 questions/month = 1,000 questions = 3,000 API calls/month',
        Inches(0.4), ct + Inches(0.05), Inches(12.5), Inches(0.38),
        size=11, italic=True, color=GRAY)
    ppt_table(s,
        ['Cost Item', 'Monthly', 'Annual', 'Notes'],
        [
            ['Claude API',           '~$35 – $50',        '~$420 – $600',    'Only external cost — scales with usage'],
            ['Hosting',              '$0',                 '$0',               'Internal BA subdomain'],
            ['Analytics Database',   '$0',                 '$0',               'Existing KpiReports SQL Server'],
            ['Application Database', '$0',                 '$0',               'New DB on existing SQL Server'],
            ['Maintenance',          'Internal',           'Internal',         'Owned by Analytics team'],
            ['TOTAL',                '~$35–$50/month',     '~$420–$600/year',  'Highly cost-effective for an enterprise tool'],
        ],
        Inches(0.7), ct + Inches(0.55), Inches(11.9), Inches(4.3),
        highlight_last_row=True)

    rect(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.85),
         PALE_GREEN, BA_GREEN, Pt(1))
    txb(s, 'The only recurring external cost is the Anthropic API  —  under $50/month for a 10-user deployment',
        Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.7),
        size=12, bold=True, color=RGB(0x00, 0x6b, 0x45), align=PP_ALIGN.CENTER)

    # ── S10: RECOMMENDATION ───────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Recommendation — Scenario 3: Hybrid Model')
    txb(s, 'Analytics team leads  ·  Technology team provides infrastructure support  ·  Claude AI accelerates delivery',
        Inches(0.4), ct + Inches(0.08), Inches(12.5), Inches(0.4),
        size=12, italic=True, color=BA_TEAL)

    reasons = [
        ('Domain Ownership',
         'Analytics team understands KPIs and data rules. Best placed to build and maintain long-term.',
         BA_BLUE),
        ('Claude AI Lowers the Bar',
         'Python, FastAPI, React — achievable by analysts with Claude Code assistance.',
         BA_NAVY),
        ('Bounded Tech Involvement',
         'Tech role: hosting, DB setup, deployment. Defined scope, no ongoing dependency.',
         BA_BLUE),
        ('Long-Term Sustainability',
         'Tool ownership stays with the team closest to the data from day one.',
         BA_NAVY),
        ('Best Value',
         'Fewer licences than Scenario 2. Less Tech dependency than Scenario 1.',
         BA_GREEN),
    ]
    for i, (title, desc, color) in enumerate(reasons):
        col = i % 2 if i < 4 else 0
        row = i // 2
        lft = Inches(0.4) + col * Inches(6.45) if i < 4 else Inches(3.7)
        tp  = ct + Inches(0.6) + row * Inches(1.82)
        w   = Inches(12.5) if i == 4 else Inches(6.1)
        rect(s, lft, tp, w, Inches(1.68), LIGHT, color, Pt(1))
        txb(s, title, lft + Inches(0.2), tp + Inches(0.15),
            w - Inches(0.4), Inches(0.45), size=12, bold=True, color=color)
        txb(s, desc, lft + Inches(0.2), tp + Inches(0.65),
            w - Inches(0.4), Inches(0.9), size=10.5, color=GRAY, wrap=True)

    # ── S11: NEXT STEPS ───────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Next Steps')

    steps = [
        ('1', 'Approve Proposal',     'Review, align on scope and scenario',                         BA_BLUE),
        ('2', 'Select Scenario',      'Confirm team structure — hybrid model recommended',            BA_NAVY),
        ('3', 'Identify Team',        'Name Analytics leads and Technology support contacts',        BA_BLUE),
        ('4', 'Procure Licences',     'Additional Claude development licences for the team',         BA_NAVY),
        ('5', 'Setup App Database',   'User management, access control, query history storage',      BA_BLUE),
        ('6', 'Configure Hosting',    'Set up paneliq.borderlessaccess.com subdomain',               BA_NAVY),
        ('7', 'Phase 1 Deployment',   'POC is ~80% ready — faster than a fresh build',               BA_TEAL),
        ('8', 'Define Access Matrix', 'Map users to data domains, roles, and field permissions',     BA_TEAL),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        col = i % 2; row = i // 2
        lft = Inches(0.4) + col * Inches(6.45)
        tp  = ct + Inches(0.22) + row * Inches(1.45)
        c = oval(s, lft, tp + Inches(0.28), Inches(0.65), Inches(0.65), color)
        ctf = c.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(14); cr.font.bold = True
        cr.font.color.rgb = WHITE; cr.font.name = 'Arial'
        tb_s = s.shapes.add_textbox(lft + Inches(0.8), tp + Inches(0.12),
                                     Inches(5.4), Inches(1.15))
        tf = tb_s.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = title
        r1.font.bold = True; r1.font.size = Pt(12)
        r1.font.color.rgb = color; r1.font.name = 'Arial'
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = desc
        r2.font.size = Pt(10.5); r2.font.color.rgb = GRAY; r2.font.name = 'Arial'

    # ── S12: THANK YOU ────────────────────────────────────────────────────────
    s = blank(prs)
    rect(s, 0, 0, Inches(8.7), prs.slide_height, BA_NAVY)
    rect(s, 0, Inches(3.5), Inches(8.7), Pt(4), BA_TEAL)
    rect(s, 0, 0, Inches(0.18), prs.slide_height, BA_BLUE)
    rect(s, Inches(8.7), 0, Inches(4.63), prs.slide_height, WHITE)

    txb(s, 'Thank You', Inches(0.5), Inches(1.6), Inches(7.9), Inches(1.4),
        size=52, bold=True, color=WHITE)
    txb(s, 'Questions & Discussion',
        Inches(0.5), Inches(3.1), Inches(7.9), Inches(0.65),
        size=20, color=BA_TEAL)
    txb(s, 'paneliq.borderlessaccess.com  |  Borderless Access  |  May 2026',
        Inches(0.5), Inches(4.4), Inches(7.9), Inches(0.55),
        size=13, color=MID_GRAY)

    prs.save(path)
    print(f'  Detailed deck saved: {path}')


# ══════════════════════════════════════════════════════════════════════════════
#  EXECUTIVE DECK  (5 slides)
# ══════════════════════════════════════════════════════════════════════════════

def exec_ppt(path):
    prs = load_template()

    # ── S1: HOOK ──────────────────────────────────────────────────────────────
    s = blank(prs)
    rect(s, 0, 0, Inches(8.7), prs.slide_height, BA_NAVY)
    rect(s, 0, Inches(4.2), Inches(8.7), Pt(4), BA_TEAL)
    rect(s, 0, 0, Inches(0.18), prs.slide_height, BA_BLUE)
    rect(s, Inches(8.7), 0, Inches(4.63), prs.slide_height, WHITE)

    # POC badge on white side
    rect(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.65), BA_GREEN)
    txb(s, '✓  POC Built — Phase 1 is 80% Complete',
        Inches(9.05), Inches(1.55), Inches(3.9), Inches(0.55),
        size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txb(s, 'What if your leadership team',
        Inches(0.4), Inches(1.2), Inches(7.9), Inches(0.75),
        size=28, color=MID_GRAY)
    txb(s, 'could get data answers\ninstantly —',
        Inches(0.4), Inches(1.9), Inches(7.9), Inches(1.4),
        size=36, bold=True, color=WHITE)
    txb(s, 'instead of waiting 1–3 days?',
        Inches(0.4), Inches(3.2), Inches(7.9), Inches(0.75),
        size=32, bold=True, color=BA_TEAL)

    txb(s, 'PanelIQ is an AI-powered analytics assistant that gives '
           'BA leadership instant, self-service access to panel '
           'performance data — in plain English.',
        Inches(0.4), Inches(4.35), Inches(7.9), Inches(1.15),
        size=13, color=MID_GRAY, wrap=True)

    txb(s, 'PanelIQ  |  Borderless Access  |  May 2026',
        Inches(0.4), Inches(6.85), Inches(7.0), Inches(0.45),
        size=10, color=RGB(0x55, 0x70, 0x90))

    # ── S2: THE PROBLEM ───────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'The Challenge — How It Works Today')

    rect(s, Inches(0.3), ct + Inches(0.1), Inches(6.0), Inches(5.85),
         RGB(0xff, 0xf5, 0xf5), RED_SOFT, Pt(1))
    txb(s, '😟  Today',
        Inches(0.5), ct + Inches(0.25), Inches(5.6), Inches(0.5),
        size=16, bold=True, color=RED_SOFT)
    pain = [
        ('Static dashboards only',
         'Tableau & SSRS show pre-defined views. Dynamic = manual request.'),
        ('1–3 business days',
         'Every ad-hoc question routed to Technology, Finance, or Analytics teams.'),
        ('Once a month or quarter',
         'Leadership receives dynamic insights only in scheduled business reviews.'),
        ('No self-service',
         'Non-technical users have no independent way to access live data.'),
    ]
    for i, (b, d) in enumerate(pain):
        two_run(s, Inches(0.5), ct + Inches(0.88) + i * Inches(1.18),
                Inches(5.7), Inches(1.05),
                '✗  ' + b + '\n', 11, True, RED_SOFT,
                d, 10.5, False, GRAY)

    rect(s, Inches(6.55), ct + Inches(0.1), Pt(3), Inches(5.85), BA_TEAL)

    rect(s, Inches(6.8), ct + Inches(0.1), Inches(6.1), Inches(5.85),
         PALE_GREEN, BA_GREEN, Pt(1))
    txb(s, '😊  With PanelIQ',
        Inches(7.0), ct + Inches(0.25), Inches(5.7), Inches(0.5),
        size=16, bold=True, color=BA_GREEN)
    good = [
        ('Ask in plain English',
         'Type any question — get a data table, chart, and analyst summary in seconds.'),
        ('Seconds, not days',
         'No waiting. No tickets. No dependency on other teams.'),
        ('Always available',
         'Leadership can query data any time — daily if needed, not just in reviews.'),
        ('Secure self-service',
         'Role-based access — each user sees only what they are permitted to see.'),
    ]
    for i, (b, d) in enumerate(good):
        two_run(s, Inches(7.0), ct + Inches(0.88) + i * Inches(1.18),
                Inches(5.7), Inches(1.05),
                '✓  ' + b + '\n', 11, True, BA_GREEN,
                d, 10.5, False, GRAY)

    txb(s, '→', Inches(6.2), ct + Inches(2.7), Inches(0.4), Inches(0.6),
        size=28, bold=True, color=BA_TEAL, align=PP_ALIGN.CENTER)

    # ── S3: SOLUTION ──────────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'PanelIQ — How It Works')

    txb(s, 'A secure web-based AI assistant at paneliq.borderlessaccess.com — '
           'ask in plain English, get answers in seconds.',
        Inches(0.4), ct + Inches(0.1), Inches(12.5), Inches(0.42),
        size=11.5, italic=True, color=GRAY)

    steps = [
        ('1', 'Ask',           '"Show me completes by\nmarket for Q1 2026"',    BA_BLUE),
        ('2', 'Translate',     'Claude AI converts to\na secure SQL query',      BA_NAVY),
        ('3', 'Query',         'Runs against internal\nKpiReports database',     BA_TEAL),
        ('4', 'Answer',        'Chart + table + analyst\nsummary in seconds',    BA_GREEN),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        lft = Inches(0.35) + i * Inches(3.25)
        tp  = ct + Inches(0.62)
        rect(s, lft, tp, Inches(3.05), Inches(2.85), LIGHT, color, Pt(1))
        c = oval(s, lft + Inches(1.0), tp + Inches(0.2), Inches(1.05), Inches(1.05), color)
        ctf = c.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(22); cr.font.bold = True
        cr.font.color.rgb = WHITE; cr.font.name = 'Arial'
        txb(s, title, lft + Inches(0.1), tp + Inches(1.38),
            Inches(2.85), Inches(0.45), size=12, bold=True, color=color,
            align=PP_ALIGN.CENTER)
        txb(s, desc, lft + Inches(0.1), tp + Inches(1.9),
            Inches(2.85), Inches(0.85), size=10.5, color=GRAY,
            align=PP_ALIGN.CENTER, wrap=True)
        if i < 3:
            txb(s, '→', Inches(3.35) + i * Inches(3.25), ct + Inches(1.7),
                Inches(0.3), Inches(0.5), size=22, bold=True, color=BA_TEAL)

    # Sample questions
    rect(s, Inches(0.35), Inches(4.05), Inches(12.6), Inches(2.85), BA_NAVY)
    txb(s, 'What your team can ask — examples',
        Inches(0.55), Inches(4.18), Inches(12.0), Inches(0.4),
        size=12, bold=True, color=BA_TEAL)
    questions = [
        '"Show me completes by market for Q1 2026"',
        '"Which vendors have the highest screenout rate?"',
        '"What is the click rate and conversion rate by country?"',
        '"Compare monthly completes trend for last 12 months"',
        '"Which projects are live right now?"',
        '"Give me overall panel health summary for this year"',
    ]
    for i, q in enumerate(questions):
        col = i % 3; row = i // 3
        lft = Inches(0.55) + col * Inches(4.15)
        tp  = Inches(4.7) + row * Inches(0.95)
        rect(s, lft, tp, Inches(3.95), Inches(0.8),
             RGB(0x22, 0x44, 0x66), BA_TEAL, Pt(1))
        txb(s, q, lft + Inches(0.12), tp + Inches(0.12),
            Inches(3.7), Inches(0.6), size=10, color=WHITE, wrap=True)

    # ── S4: INVESTMENT ────────────────────────────────────────────────────────
    s = blank(prs)
    ct = slide_header(s, prs, 'Investment & Implementation')

    stats = [
        ('~$50/mo',  'Maximum ongoing\noperating cost\n(API only)',           PALE_BLUE,  BA_BLUE,  GRAY),
        ('80%',      'Phase 1 already built\n(POC validated —\nfaster go-live)', PALE_GREEN, BA_GREEN, GRAY),
        ('Seconds',  'Time to answer any\npanel question\nvs 1–3 days today',  LIGHT,      BA_TEAL,  GRAY),
    ]
    for i, (big, label, bg, bc, lc) in enumerate(stats):
        stat_block(s, Inches(0.4) + i * Inches(4.3), Inches(1.1),
                   Inches(4.05), Inches(2.1), big, label, bg, bc, lc)

    txb(s, 'Implementation Options',
        Inches(0.4), Inches(3.35), Inches(12.5), Inches(0.38),
        size=13, bold=True, color=BA_NAVY)

    scenarios = [
        ('Scenario 1', 'Technology Team',       '~6–8 months',   '~$1,200–$1,600\ndev licences', LIGHT,      BA_BLUE),
        ('Scenario 2', 'Analytics Team',        '~12–15 months', '~$2,400–$4,500\ndev licences', LIGHT,      BA_NAVY),
        ('Scenario 3', 'Hybrid  ★ Recommended', '~8–10 months',  '~$1,600–$2,000\ndev licences', PALE_GREEN, BA_GREEN),
    ]
    for i, (sc, team, timeline, cost, bg, color) in enumerate(scenarios):
        lft = Inches(0.4) + i * Inches(4.3)
        tp  = Inches(3.82)
        rect(s, lft, tp, Inches(4.05), Inches(3.0), bg, color, Pt(1))
        txb(s, sc,    lft + Inches(0.15), tp + Inches(0.15),
            Inches(3.75), Inches(0.32), size=10, bold=True, color=MID_GRAY)
        txb(s, team,  lft + Inches(0.15), tp + Inches(0.45),
            Inches(3.75), Inches(0.45), size=13, bold=True, color=color)
        txb(s, f'⏱  Timeline:  {timeline}',
            lft + Inches(0.15), tp + Inches(1.05),
            Inches(3.75), Inches(0.42), size=11, color=GRAY)
        txb(s, f'💰  Dev licences:  {cost}',
            lft + Inches(0.15), tp + Inches(1.52),
            Inches(3.75), Inches(0.68), size=11, color=GRAY, wrap=True)
        txb(s, 'Ongoing API: ~$35–$50/month',
            lft + Inches(0.15), tp + Inches(2.55),
            Inches(3.75), Inches(0.35), size=9.5, italic=True, color=color)

    # ── S5: THE ASK ───────────────────────────────────────────────────────────
    s = blank(prs)
    rect(s, 0, 0, Inches(8.7), prs.slide_height, BA_NAVY)
    rect(s, 0, Inches(0.5), Inches(8.7), Pt(4), BA_TEAL)
    rect(s, 0, 0, Inches(0.18), prs.slide_height, BA_BLUE)
    rect(s, Inches(8.7), 0, Inches(4.63), prs.slide_height, WHITE)

    txb(s, 'One Decision. Three Steps. Live in Weeks.',
        Inches(0.4), Inches(0.15), Inches(8.1), Inches(0.7),
        size=24, bold=True, color=WHITE)

    actions = [
        ('01', 'Approve',
         'Confirm scope and team model\n(Hybrid recommended)',
         BA_BLUE),
        ('02', 'Assign',
         'Name Analytics leads + Tech support\nProcure Claude licences',
         BA_TEAL),
        ('03', 'Launch',
         'Phase 1 goes live at\npaneliq.borderlessaccess.com\nwithin 5–7 weeks',
         BA_GREEN),
    ]
    for i, (num, title, desc, color) in enumerate(actions):
        tp = Inches(1.1) + i * Inches(1.75)
        rect(s, Inches(0.35), tp, Inches(8.1), Inches(1.6), color)
        c = oval(s, Inches(0.45), tp + Inches(0.28),
                 Inches(1.05), Inches(1.05), WHITE)
        ctf = c.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(20); cr.font.bold = True
        cr.font.color.rgb = color; cr.font.name = 'Arial'
        txb(s, title, Inches(1.65), tp + Inches(0.2),
            Inches(6.6), Inches(0.5), size=15, bold=True, color=WHITE)
        txb(s, desc, Inches(1.65), tp + Inches(0.75),
            Inches(6.6), Inches(0.75), size=11, color=WHITE, wrap=True)

    txb(s, 'paneliq.borderlessaccess.com  |  Borderless Access  |  May 2026',
        Inches(0.35), Inches(6.6), Inches(8.1), Inches(0.45),
        size=10, color=MID_GRAY)

    # Right side summary on white panel
    txb(s, 'Why now?',
        Inches(9.0), Inches(1.5), Inches(4.0), Inches(0.45),
        size=14, bold=True, color=BA_BLUE)
    why = [
        '✓  POC is 80% built',
        '✓  Under $50/month to run',
        '✓  No new infrastructure needed',
        '✓  Analytics team can lead it',
        '✓  Claude AI accelerates delivery',
    ]
    for i, w in enumerate(why):
        txb(s, w, Inches(9.0), Inches(2.1) + i * Inches(0.72),
            Inches(4.0), Inches(0.6), size=12, color=BA_NAVY)

    prs.save(path)
    print(f'  Executive deck saved: {path}')


# ── MAIN ──────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    detailed_ppt(os.path.join(BASE, 'PanelIQ_Proposal_May2026.pptx'))
    exec_ppt(os.path.join(BASE,     'PanelIQ_Executive_Deck_May2026.pptx'))
    print('Done.')
